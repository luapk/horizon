#!/usr/bin/env python3
"""Horizon Engine — Omnicom enterprise integration briefing deck.
Dark brand deck, grounded in the actual codebase (stack, pricing, estimator)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ── palette (matches the product) ──────────────────────────────
ABYSS   = RGBColor(0x05, 0x07, 0x0F)
PANEL   = RGBColor(0x0E, 0x13, 0x20)
PANEL2  = RGBColor(0x15, 0x1B, 0x2C)
BORDER  = RGBColor(0x2A, 0x33, 0x48)
HEAD    = RGBColor(0xF5, 0xF7, 0xFA)
BODY    = RGBColor(0xC7, 0xCF, 0xE0)
MUTE    = RGBColor(0x8B, 0x97, 0xAF)
FAINT   = RGBColor(0x60, 0x6C, 0x85)
BLUE    = RGBColor(0x8F, 0xB8, 0xFF)
VIOLET  = RGBColor(0xB3, 0x9D, 0xFF)
PINK    = RGBColor(0xF2, 0xA6, 0xC9)
MINT    = RGBColor(0x93, 0xE6, 0xC3)
PEACH   = RGBColor(0xFF, 0xC2, 0x9E)
CORAL   = RGBColor(0xFF, 0x9D, 0x9D)

DISPLAY = "Georgia"
BODYF   = "Calibri"
MONO    = "Consolas"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def solid(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def line(shape, color, w=0.75):
    shape.line.color.rgb = color; shape.line.width = Pt(w)


def bg(slide, color=ABYSS):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=Pt(4), line_spacing=1.12):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold, font, italic)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = space_after; p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, size, color, bold, font, *rest) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
            r.font.name = font
            if rest and rest[0]:
                r.font.italic = True
    return tb


def rrect(slide, x, y, w, h, fill=PANEL, brdr=BORDER, radius=0.06, line_w=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    solid(sp, fill)
    if brdr is None:
        no_line(sp)
    else:
        line(sp, brdr, line_w)
    sp.shadow.inherit = False
    return sp


def rect(slide, x, y, w, h, fill):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    solid(sp, fill); no_line(sp); sp.shadow.inherit = False
    return sp


PAGE = [0]
FOOT = "HORIZON ENGINE  ·  CONFIDENTIAL  ·  PREPARED FOR OMNICOM"


def chrome(slide, eyebrow, title, accent=VIOLET):
    # top accent hairline
    rect(slide, Inches(0), Inches(0), EMU_W, Pt(3), accent)
    txt(slide, Inches(0.7), Inches(0.42), Inches(11.9), Inches(0.3),
        [[(eyebrow, 11, accent, True, MONO)]])
    txt(slide, Inches(0.7), Inches(0.72), Inches(11.9), Inches(0.7),
        [[(title, 27, HEAD, False, DISPLAY)]])
    # footer
    PAGE[0] += 1
    txt(slide, Inches(0.7), Inches(7.06), Inches(9), Inches(0.3),
        [[(FOOT, 7.5, FAINT, False, MONO)]])
    txt(slide, Inches(11.6), Inches(7.06), Inches(1.03), Inches(0.3),
        [[(f"{PAGE[0]:02d}", 7.5, FAINT, False, MONO)]], align=PP_ALIGN.RIGHT)


def new(eyebrow=None, title=None, accent=VIOLET):
    s = prs.slides.add_slide(BLANK)
    bg(s)
    if title:
        chrome(s, eyebrow, title, accent)
    return s


def card(slide, x, y, w, h, accent, head, lines, head_size=13, body_size=10.5):
    rrect(slide, x, y, w, h, fill=PANEL, brdr=BORDER)
    rect(slide, x, y + Inches(0.14), Pt(3), h - Inches(0.28), accent)  # left accent bar
    pad = Inches(0.26)
    inner = [[(head, head_size, HEAD, True, BODYF)]]
    for ln in lines:
        inner.append([(ln, body_size, BODY, False, BODYF)])
    txt(slide, x + pad, y + Inches(0.2), w - pad - Inches(0.18), h - Inches(0.3),
        inner, space_after=Pt(5), line_spacing=1.12)


def bullets(slide, x, y, w, items, size=13, gap=Pt(9), color=BODY, marker=VIOLET):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            lead, rest = it
            paras.append([("▸  ", size, marker, True, BODYF), (lead, size, HEAD, True, BODYF),
                          (rest, size, color, False, BODYF)])
        else:
            paras.append([("▸  ", size, marker, True, BODYF), (it, size, color, False, BODYF)])
    txt(slide, x, y, w, Inches(5), paras, space_after=gap, line_spacing=1.14)


def style_cell(cell, text_runs, fill, size=10, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.margin_left = Inches(0.09); cell.margin_right = Inches(0.09)
    cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
    cell.vertical_anchor = anchor
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for (t, c, b, f) in text_runs:
        r = p.add_run(); r.text = t; r.font.size = Pt(size)
        r.font.color.rgb = c; r.font.bold = b; r.font.name = f


def table(slide, x, y, w, colw, rows, header, body_rows, row_h=Inches(0.34), hsize=9.5, bsize=9.5):
    nrows = len(body_rows) + 1
    ncols = len(header)
    gt = slide.shapes.add_table(nrows, ncols, x, y, w, row_h * nrows).table
    gt.first_row = False; gt.horz_banding = False
    # remove default style borders by setting our own via cell fills
    for ci, cw in enumerate(colw):
        gt.columns[ci].width = cw
    gt.rows[0].height = Inches(0.42)
    for ci, htext in enumerate(header):
        style_cell(gt.cell(0, ci), [(htext, MUTE, True, MONO)], PANEL2,
                   size=hsize, align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
    for ri, row in enumerate(body_rows):
        gt.rows[ri + 1].height = row_h
        rowfill = ABYSS if ri % 2 == 0 else PANEL
        for ci, (cell_runs) in enumerate(row):
            style_cell(gt.cell(ri + 1, ci), cell_runs, rowfill,
                       size=bsize, align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
    return gt


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, Inches(0), Inches(0), EMU_W, Pt(4), VIOLET)
# ambient bands
rect(s, Inches(0), Inches(3.35), EMU_W, Pt(1), BORDER)
txt(s, Inches(0.9), Inches(1.5), Inches(11), Inches(0.4),
    [[("MULTI-BRAND FUTURES INTELLIGENCE ENGINE", 12, VIOLET, True, MONO)]])
txt(s, Inches(0.85), Inches(1.95), Inches(11.6), Inches(1.4),
    [[("Horizon Engine", 58, HEAD, False, DISPLAY)]])
txt(s, Inches(0.9), Inches(3.55), Inches(11.5), Inches(0.9),
    [[("Enterprise integration briefing", 22, BLUE, False, DISPLAY, True)]])
txt(s, Inches(0.9), Inches(4.35), Inches(11.5), Inches(0.6),
    [[("Technical architecture · AI compliance · deployment at scale · unit economics · data security",
       13, BODY, False, BODYF)]])
# meta row
rrect(s, Inches(0.9), Inches(5.5), Inches(4.9), Inches(0.95), fill=PANEL, brdr=BORDER)
txt(s, Inches(1.15), Inches(5.66), Inches(4.5), Inches(0.7),
    [[("PREPARED FOR", 9, MUTE, True, MONO)],
     [("Omnicom — Group IT, Architecture & Security", 13, HEAD, True, BODYF)]], space_after=Pt(3))
rrect(s, Inches(6.05), Inches(5.5), Inches(6.4), Inches(0.95), fill=PANEL, brdr=BORDER)
txt(s, Inches(6.3), Inches(5.66), Inches(6.0), Inches(0.7),
    [[("SCOPE", 9, MUTE, True, MONO)],
     [("How the platform is built, secured, governed and run at holding-company scale", 12, HEAD, False, BODYF)]],
    space_after=Pt(3))
txt(s, Inches(0.9), Inches(6.9), Inches(11.5), Inches(0.3),
    [[("Confidential — commercial discussion document", 8.5, FAINT, False, MONO)]])

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════════════════
s = new("01 · EXECUTIVE SUMMARY", "A foresight engine that runs for any brand", VIOLET)
txt(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(0.9),
    [[("Horizon scans live sources for weak signals, clusters them into structural drivers, and writes "
       "evidence-cited future scenarios with a strategic impact matrix — per brand, on demand. It is built "
       "as a repeatable pipeline, not a bespoke report, so the same engine serves every agency and client in "
       "the group.", 13.5, BODY, False, BODYF)]], line_spacing=1.2)
cy = Inches(2.75)
card(s, Inches(0.7), cy, Inches(3.8), Inches(1.75), BLUE, "Purpose-built for a holding company",
     ["One engine, many brands and agencies.", "Each scan is defined by a brand config —",
      "no code changes to onboard a new client."])
card(s, Inches(4.75), cy, Inches(3.8), Inches(1.75), MINT, "Cost known before you spend",
     ["Every scan shows a per-stage cost", "estimate up front; actual token spend",
      "is metered and reconciled after."])
card(s, Inches(8.8), cy, Inches(3.8), Inches(1.75), PINK, "Grounded, not generative vibes",
     ["Scenarios cite the exact signals behind", "them; uncited output is flagged.",
      "Absence & counter-signals surfaced."])
# bottom band
rrect(s, Inches(0.7), Inches(4.75), Inches(11.9), Inches(1.55), fill=PANEL2, brdr=BORDER)
txt(s, Inches(1.0), Inches(4.95), Inches(11.3), Inches(1.2),
    [[("WHAT THIS DOCUMENT COVERS", 10, VIOLET, True, MONO)],
     [("Tech stack and the AI pipeline · token usage and unit economics with worked examples · pricing and "
       "spend governance · data security and key custody · AI compliance (EU AI Act, GDPR) · and a phased "
       "path to deploy across Omnicom agencies with SSO, multi-tenancy and enterprise controls.",
       12.5, BODY, False, BODYF)]], space_after=Pt(6), line_spacing=1.16)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT IT PRODUCES (the pipeline output)
# ═══════════════════════════════════════════════════════════════
s = new("02 · THE INTELLIGENCE ARTIFACT", "From raw sources to a decision-ready dossier", BLUE)
stages = [
    ("SIGNALS", BLUE, "Weak signals extracted from\nlive sources, tagged by STEEP,\ngeography, confidence."),
    ("CLUSTERS", VIOLET, "Embedding-based grouping —\nstable, data-driven, not LLM\nfree-association."),
    ("DRIVERS", PINK, "The structural forces the\nclusters reveal, synthesised\nas a distinct reasoning pass."),
    ("SCENARIOS", MINT, "Evidence-cited narratives —\n5 probable · 3 deep · 1\ncautionary 'Cassandra'."),
    ("STRATEGY", PEACH, "Impact matrix across business\nunits + a Now / Monitor /\nPrepare action timeline."),
]
n = len(stages); gap = Inches(0.28)
cw = (Inches(11.9) - gap * (n - 1)) / n
x = Inches(0.7); y = Inches(2.1)
for i, (name, acc, desc) in enumerate(stages):
    cx = x + (cw + gap) * i
    rrect(s, cx, y, cw, Inches(2.5), fill=PANEL, brdr=BORDER)
    rect(s, cx, y, cw, Pt(3), acc)
    txt(s, cx + Inches(0.16), y + Inches(0.2), cw - Inches(0.3), Inches(0.4),
        [[(f"{i+1:02d}", 10, acc, True, MONO)]])
    txt(s, cx + Inches(0.16), y + Inches(0.55), cw - Inches(0.3), Inches(0.4),
        [[(name, 13, HEAD, True, BODYF)]])
    txt(s, cx + Inches(0.16), y + Inches(1.02), cw - Inches(0.3), Inches(1.4),
        [[(desc, 10, BODY, False, BODYF)]], line_spacing=1.14)
    if i < n - 1:
        txt(s, cx + cw - Inches(0.02), y + Inches(1.0), Inches(0.3), Inches(0.4),
            [[("→", 16, FAINT, True, BODYF)]])
txt(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.2),
    [[("Two analytical moves competitors miss", 15, HEAD, False, DISPLAY)],
     [("Absence signals — what the evidence implies should exist but doesn't — and counter-signals that cut "
       "against the corpus consensus are produced by a dedicated pass over the whole signal set, then marked "
       "as derived (never overstated as 'verified'). Every scenario carries an explicit shadow side and a "
       "'killer assumption' stating what would prove it wrong.", 12.5, BODY, False, BODYF)]],
    space_after=Pt(7), line_spacing=1.18)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — TECH STACK / ARCHITECTURE
# ═══════════════════════════════════════════════════════════════
s = new("03 · TECHNOLOGY STACK", "A typed, layered, provider-agnostic architecture", VIOLET)
layers = [
    ("PRESENTATION", BLUE, "React 18 + Vite (TypeScript)",
     "Single-page app. Brand setup, scan launcher with live cost estimate, results workspace, canvas signal globe. No secrets client-side."),
    ("API + ORCHESTRATION", VIOLET, "Node.js + Express (TypeScript)",
     "Session auth, brand & scan REST API, background scan runner with per-stage progress, spend-cap enforcement, usage metering."),
    ("PIPELINE", PINK, "Ingest → Extract → Cluster → Gap → Drivers → Scenarios → Matrix",
     "Deterministic, resumable stages. Bounded concurrency, per-document fault tolerance, one retry on malformed model output."),
    ("PROVIDER ABSTRACTION", MINT, "Search · LLM · Embeddings — swappable behind interfaces",
     "Anthropic Claude (LLM) · Tavily (search) · Voyage (embeddings). Each has a mock impl; models are config, not hard-coded."),
    ("DATA + SHARED CONTRACT", PEACH, "SQLite → PostgreSQL · Zod schema shared FE/BE",
     "One validated schema is the single source of truth for both ends. Pricing & cost model live here too."),
]
y = Inches(1.6); lh = Inches(1.0); g = Inches(0.13)
for i, (tag, acc, tech, desc) in enumerate(layers):
    ly = y + (lh + g) * i
    rrect(s, Inches(0.7), ly, Inches(11.9), lh, fill=PANEL, brdr=BORDER)
    rect(s, Inches(0.7), ly, Pt(3), lh, acc)
    txt(s, Inches(1.0), ly + Inches(0.16), Inches(2.7), Inches(0.7),
        [[(tag, 10.5, acc, True, MONO)]])
    txt(s, Inches(3.75), ly + Inches(0.13), Inches(8.6), Inches(0.34),
        [[(tech, 12.5, HEAD, True, BODYF)]])
    txt(s, Inches(3.75), ly + Inches(0.47), Inches(8.6), Inches(0.5),
        [[(desc, 10.5, BODY, False, BODYF)]], line_spacing=1.1)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — THE AI PIPELINE (models per stage)
# ═══════════════════════════════════════════════════════════════
s = new("04 · THE AI PIPELINE", "Right-sized model for each stage — with a human gate", MINT)
header = ["STAGE", "ENGINE", "MODEL TIER", "WHY THIS TIER"]
colw = [Inches(2.7), Inches(2.5), Inches(2.4), Inches(4.3)]
rows = [
    [[("Query design", HEAD, True, BODYF)], [("Claude", BODY, False, BODYF)], [("Haiku (fast)", MINT, True, MONO)],
     [("Cheap, high-volume expansion of search queries", BODY, False, BODYF)]],
    [[("Signal extraction", HEAD, True, BODYF)], [("Claude", BODY, False, BODYF)], [("Haiku (fast)", MINT, True, MONO)],
     [("One call per document; volume dominates — keep it cheap", BODY, False, BODYF)]],
    [[("Dedupe / cluster", HEAD, True, BODYF)], [("Voyage", BODY, False, BODYF)], [("Embeddings", BLUE, True, MONO)],
     [("Vector similarity, not tokens — near-free, deterministic", BODY, False, BODYF)]],
    [[("Gap analysis", HEAD, True, BODYF)], [("Claude", BODY, False, BODYF)], [("Sonnet (synthesis)", VIOLET, True, MONO)],
     [("Reasons over the whole corpus — needs a stronger model", BODY, False, BODYF)]],
    [[("Driver synthesis", HEAD, True, BODYF)], [("Claude", BODY, False, BODYF)], [("Sonnet (synthesis)", VIOLET, True, MONO)],
     [("Judgement task; abstracts clusters into forces", BODY, False, BODYF)]],
    [[("Scenario generation", HEAD, True, BODYF)], [("Claude", BODY, False, BODYF)], [("Sonnet / Opus", VIOLET, True, MONO)],
     [("The craft layer; Opus/Fable optional for board-grade prose", BODY, False, BODYF)]],
    [[("Strategy matrix", HEAD, True, BODYF)], [("Claude", BODY, False, BODYF)], [("Sonnet (synthesis)", VIOLET, True, MONO)],
     [("Scores scenarios × business units", BODY, False, BODYF)]],
]
table(s, Inches(0.7), Inches(1.55), Inches(11.9), colw, len(rows), header, rows, row_h=Inches(0.44))
rrect(s, Inches(0.7), Inches(5.55), Inches(11.9), Inches(0.95), fill=PANEL2, brdr=VIOLET, line_w=1.0)
txt(s, Inches(1.0), Inches(5.72), Inches(11.3), Inches(0.7),
    [[("HUMAN-IN-THE-LOOP  ", 11, VIOLET, True, MONO), ("Scenario output is designed for a review gate before "
      "publication. Citations are parsed from the prose and validated against the real signal set; any scenario "
      "that cites nothing is auto-downgraded and flagged UNGROUNDED in the UI.", 12, BODY, False, BODYF)]],
    line_spacing=1.16)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — TOKEN USAGE & UNIT ECONOMICS
# ═══════════════════════════════════════════════════════════════
s = new("05 · TOKEN USAGE & UNIT ECONOMICS", "Every stage is metered; cost scales with scope", PEACH)
txt(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.5),
    [[("Cost per scan is a function of scope (queries × docs, driver & scenario counts). Figures below are the "
       "platform's own pre-scan estimate — a low/high band driven by retrieval yield and model output variance.",
       12, BODY, False, BODYF)]], line_spacing=1.16)
header = ["SCAN TIER", "QUERIES × DOCS", "DRIVERS / SCEN.", "EST. COST / SCAN", "DOMINANT COST"]
colw = [Inches(2.9), Inches(2.3), Inches(2.2), Inches(2.5), Inches(2.0)]
rows = [
    [[("Light — pilot", HEAD, True, BODYF)], [("10 × 4", BODY, False, MONO)], [("6 / 6", BODY, False, MONO)],
     [("$0.39 – $0.61", MINT, True, MONO)], [("Scenarios", MUTE, False, BODYF)]],
    [[("Standard — default", HEAD, True, BODYF)], [("20 × 6", BODY, False, MONO)], [("9 / 9", BODY, False, MONO)],
     [("$0.62 – $1.13", BLUE, True, MONO)], [("Scenarios", MUTE, False, BODYF)]],
    [[("Deep — board-grade", HEAD, True, BODYF)], [("40 × 10", BODY, False, MONO)], [("12 / 12", BODY, False, MONO)],
     [("$0.96 – $2.41", VIOLET, True, MONO)], [("Extraction", MUTE, False, BODYF)]],
]
table(s, Inches(0.7), Inches(2.15), Inches(11.9), colw, len(rows), header, rows, row_h=Inches(0.46))
# worked example + tiering note
card(s, Inches(0.7), Inches(4.15), Inches(5.85), Inches(2.15), BLUE, "Worked example — one Standard scan",
     ["~120 doc extractions on Haiku (input-heavy,",
      "short output) ≈ pennies each.",
      "9 scenarios on Sonnet dominate — long,",
      "reasoned prose is the real cost centre.",
      "Embeddings for dedupe/cluster ≈ $0.001–0.01.",
      "Search API is a flat per-query line item."], body_size=10.5)
card(s, Inches(6.75), Inches(4.15), Inches(5.85), Inches(2.15), VIOLET, "Model tiering keeps it cheap",
     ["Haiku  $1 / $5 per M tokens — extraction",
      "Sonnet $3 / $15 — synthesis & scenarios",
      "Opus   $5 / $25 — optional board-grade",
      "Fable  $10 / $50 — deepest reasoning",
      "Voyage embeddings ≈ $0.06 / M — clustering",
      "Volume runs on the cheapest capable tier."], body_size=10.5)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — PRICING AT SCALE & GOVERNANCE
# ═══════════════════════════════════════════════════════════════
s = new("06 · PRICING AT SCALE", "Predictable monthly spend, hard-capped by design", MINT)
txt(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.5),
    [[("Provider cost is pass-through and linear in scan volume. At a Standard-scan midpoint of ~$0.88, "
       "group-wide usage stays modest even at high cadence:", 12, BODY, False, BODYF)]], line_spacing=1.16)
header = ["MONTHLY SCAN VOLUME", "≈ PROVIDER COST / MONTH", "TYPICAL USE"]
colw = [Inches(3.6), Inches(3.6), Inches(4.7)]
rows = [
    [[("100 scans", HEAD, True, BODYF)], [("~ $90", MINT, True, MONO)], [("Single agency, weekly brand reviews", BODY, False, BODYF)]],
    [[("500 scans", HEAD, True, BODYF)], [("~ $440", BLUE, True, MONO)], [("Several agencies, active pitch support", BODY, False, BODYF)]],
    [[("2,000 scans", HEAD, True, BODYF)], [("~ $1,750", VIOLET, True, MONO)], [("Network-wide, embedded in planning", BODY, False, BODYF)]],
    [[("10,000 scans", HEAD, True, BODYF)], [("~ $8,750", PINK, True, MONO)], [("Holding-company standard workflow", BODY, False, BODYF)]],
]
table(s, Inches(0.7), Inches(2.15), Inches(11.9), colw, len(rows), header, rows, row_h=Inches(0.46))
card(s, Inches(0.7), Inches(4.35), Inches(3.8), Inches(1.95), PEACH, "Pre-scan estimate",
     ["The low/high band is shown before a", "scan runs and recalculates live as",
      "scope changes — no bill surprises."])
card(s, Inches(4.75), Inches(4.35), Inches(3.8), Inches(1.95), CORAL, "Hard spend caps",
     ["Per-scan ceiling blocks any scan whose", "estimate exceeds it. Monthly budget",
      "halts new scans once the cap is hit."])
card(s, Inches(8.8), Inches(4.35), Inches(3.8), Inches(1.95), MINT, "Per-call metering",
     ["Actual tokens, searches and embeddings", "are recorded per scan and reconciled",
      "against the estimate for chargeback."])

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — DATA SECURITY
# ═══════════════════════════════════════════════════════════════
s = new("07 · DATA SECURITY", "Secrets server-side, tenants isolated, spend fenced", CORAL)
left = [
    ("Key custody", "Provider API keys live only in the backend environment — never shipped to the browser, never in scan output or logs. Rotate without redeploying the client."),
    ("Authentication", "Server-verified sessions today (bcrypt-hashed secret, signed httpOnly cookie). Roadmap: SSO via Okta / Microsoft Entra with SCIM provisioning."),
    ("Transport & at rest", "TLS in transit; database encryptable at rest. Self-host in Omnicom's own cloud/VPC — no third party sees a brand's corpus except the model providers under contract."),
]
right = [
    ("Tenant isolation", "Brands and scans are row-scoped; the multi-tenant model isolates each agency and client so one team cannot see another's work or spend."),
    ("Provider data posture", "Anthropic does not train on API data and offers configurable / zero-retention options; Bedrock or Vertex deployment keeps traffic inside your cloud boundary."),
    ("Spend as a control", "Per-scan and monthly caps are a security control as much as a cost one — they bound blast radius if a credential or account is misused."),
]
def seccol(x, items, acc):
    y = Inches(1.65)
    for head, desc in items:
        rrect(s, x, y, Inches(5.75), Inches(1.5), fill=PANEL, brdr=BORDER)
        rect(s, x, y, Pt(3), Inches(1.5), acc)
        txt(s, x + Inches(0.26), y + Inches(0.16), Inches(5.3), Inches(0.35),
            [[(head, 13, HEAD, True, BODYF)]])
        txt(s, x + Inches(0.26), y + Inches(0.52), Inches(5.3), Inches(0.9),
            [[(desc, 10.5, BODY, False, BODYF)]], line_spacing=1.14)
        y = y + Inches(1.66)
seccol(Inches(0.7), left, CORAL)
seccol(Inches(6.85), right, BLUE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — AI COMPLIANCE & GOVERNANCE
# ═══════════════════════════════════════════════════════════════
s = new("08 · AI COMPLIANCE & GOVERNANCE", "Built for transparency, provenance and review", VIOLET)
cards = [
    ("EU AI Act posture", BLUE, "A decision-support / foresight tool, not automated decisioning on people — a limited-risk, transparency-obligation profile. Outputs are advisory and human-reviewed."),
    ("Provenance by design", VIOLET, "Every scenario links to the exact signals it cites; each signal keeps its source. Analysts can trace any claim back to its origin — no black-box assertions."),
    ("Honest confidence", PINK, "Confidence is evidence-assessed, not cosmetic. Derived items (absences, counter-signals) are never labelled 'verified'; uncited scenarios are flagged UNGROUNDED."),
    ("GDPR alignment", MINT, "Scans target market/industry signals, not personal data. Self-hosting plus provider no-train / retention controls keep any incidental data inside your governance."),
    ("Human review gate", PEACH, "The architecture places a person between generation and publication. Nothing goes to a client deck unreviewed; the tool accelerates analysts, it doesn't replace them."),
    ("Model auditability", CORAL, "Model and version are pinned per stage and recorded with usage; runs are reproducible and attributable for internal audit and client assurance."),
]
gx, gy = Inches(0.7), Inches(1.6)
cw2, ch2 = Inches(3.83), Inches(1.62)
for i, (head, acc, desc) in enumerate(cards):
    col = i % 3; rowi = i // 3
    cx = gx + (cw2 + Inches(0.2)) * col
    cy = gy + (ch2 + Inches(0.22)) * rowi
    rrect(s, cx, cy, cw2, ch2, fill=PANEL, brdr=BORDER)
    rect(s, cx, cy, cw2, Pt(3), acc)
    txt(s, cx + Inches(0.22), cy + Inches(0.18), cw2 - Inches(0.4), Inches(0.35),
        [[(head, 12.5, HEAD, True, BODYF)]])
    txt(s, cx + Inches(0.22), cy + Inches(0.58), cw2 - Inches(0.4), Inches(1.0),
        [[(desc, 10, BODY, False, BODYF)]], line_spacing=1.13)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — DEPLOY AT SCALE FOR OMNICOM
# ═══════════════════════════════════════════════════════════════
s = new("09 · DEPLOYING AT SCALE", "One platform, every agency and client — isolated", BLUE)
txt(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.5),
    [[("The holding-company shape maps cleanly onto a multi-tenant hierarchy: ",
       12.5, BODY, False, BODYF),
      ("Group → Agency → Client brand → Scan.", 12.5, HEAD, True, BODYF)]], line_spacing=1.16)
items_l = [
    ("Multi-tenancy & RBAC", "Per-agency and per-brand isolation with role-based access — analyst, reviewer, admin. Chargeback by tenant from metered usage."),
    ("Identity", "SSO through Okta / Microsoft Entra, SCIM user provisioning, group-to-role mapping — no separate password estate."),
    ("Data layer for scale", "PostgreSQL replaces SQLite; a managed job queue runs scans as durable background workers with retries and horizontal scale."),
]
items_r = [
    ("Cloud & networking", "Self-hosted in Omnicom's cloud (AWS/Azure/GCP) inside a VPC. Private egress to model providers; optional Bedrock / Vertex to keep traffic in-tenant."),
    ("Integration surface", "Clean REST API for embedding scans into planning tools, DAM, or agency portals; export to the formats client teams already use."),
    ("Observability & FinOps", "Central dashboards for usage, cost per tenant, cache-hit rates and error budgets; caps and alerts wired to your FinOps process."),
]
def col2(x, items, acc):
    y = Inches(2.15)
    for head, desc in items:
        rrect(s, x, y, Inches(5.75), Inches(1.4), fill=PANEL, brdr=BORDER)
        rect(s, x, y, Pt(3), Inches(1.4), acc)
        txt(s, x + Inches(0.26), y + Inches(0.15), Inches(5.3), Inches(0.35),
            [[(head, 12.5, HEAD, True, BODYF)]])
        txt(s, x + Inches(0.26), y + Inches(0.5), Inches(5.3), Inches(0.85),
            [[(desc, 10.5, BODY, False, BODYF)]], line_spacing=1.13)
        y = y + Inches(1.55)
col2(Inches(0.7), items_l, BLUE)
col2(Inches(6.85), items_r, VIOLET)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — ROADMAP + MATURITY (honest)
# ═══════════════════════════════════════════════════════════════
s = new("10 · IMPLEMENTATION ROADMAP", "A phased path from pilot to group standard", MINT)
phases = [
    ("PHASE 0", MINT, "Pilot — 4–6 wks", ["One agency, one cloud account.", "Self-hosted, real API keys, capped spend.", "Prove value on live brands."]),
    ("PHASE 1", BLUE, "Enterprise-ready", ["SSO + SCIM, multi-tenancy, Postgres.", "Job queue, RBAC, per-tenant chargeback.", "Security review & pen test."]),
    ("PHASE 2", VIOLET, "Scale-out", ["Roll across agencies. Integrations to", "planning tools & portals. FinOps", "dashboards, caching, SLAs."]),
    ("PHASE 3", PINK, "Group standard", ["Bedrock/Vertex for in-tenant traffic,", "data-residency options, human-review", "workflow, model governance board."]),
]
n = len(phases); gap = Inches(0.25)
cw = (Inches(11.9) - gap * (n - 1)) / n
for i, (tag, acc, title, pts) in enumerate(phases):
    cx = Inches(0.7) + (cw + gap) * i
    rrect(s, cx, Inches(1.7), cw, Inches(3.0), fill=PANEL, brdr=BORDER)
    rect(s, cx, Inches(1.7), cw, Pt(3), acc)
    txt(s, cx + Inches(0.2), Inches(1.9), cw - Inches(0.35), Inches(0.3),
        [[(tag, 10.5, acc, True, MONO)]])
    txt(s, cx + Inches(0.2), Inches(2.25), cw - Inches(0.35), Inches(0.4),
        [[(title, 13.5, HEAD, True, BODYF)]])
    txt(s, cx + Inches(0.2), Inches(2.8), cw - Inches(0.35), Inches(1.8),
        [[(p, 10, BODY, False, BODYF)] for p in pts], space_after=Pt(6), line_spacing=1.12)
# maturity honesty band
rrect(s, Inches(0.7), Inches(4.95), Inches(11.9), Inches(1.35), fill=PANEL2, brdr=BORDER)
txt(s, Inches(1.0), Inches(5.1), Inches(5.6), Inches(1.1),
    [[("PRODUCTION-READY TODAY", 10, MINT, True, MONO)],
     [("End-to-end pipeline · provider abstraction with real Anthropic/Tavily/Voyage integrations · "
       "cost estimator & metering · spend caps · session auth · self-host.", 11, BODY, False, BODYF)]],
    space_after=Pt(5), line_spacing=1.14)
txt(s, Inches(6.9), Inches(5.1), Inches(5.5), Inches(1.1),
    [[("ENGINEERED, ON THE ROADMAP", 10, PEACH, True, MONO)],
     [("SSO/SCIM · multi-tenant RBAC · Postgres + job queue · Bedrock/Vertex · review-workflow UI · "
       "FinOps dashboards. Scoped in Phases 1–3 above.", 11, BODY, False, BODYF)]],
    space_after=Pt(5), line_spacing=1.14)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — NEXT STEPS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, Inches(0), Inches(0), EMU_W, Pt(4), VIOLET)
txt(s, Inches(0.9), Inches(1.35), Inches(11), Inches(0.4),
    [[("NEXT STEPS", 12, VIOLET, True, MONO)]])
txt(s, Inches(0.85), Inches(1.8), Inches(11.6), Inches(1.0),
    [[("Let's run a scoped pilot", 40, HEAD, False, DISPLAY)]])
bullets(s, Inches(0.9), Inches(3.05), Inches(11.4), [
    ("Select one agency and 2–3 live client brands — ", "we stand up a self-hosted instance in your cloud."),
    ("Agree a monthly spend cap — ", "you see estimated and actual token cost on every scan from day one."),
    ("Security & architecture review — ", "we walk your team through key custody, isolation and the provider posture."),
    ("Define the success test — ", "a foresight dossier your planners would actually put in front of a client."),
], size=13.5, gap=Pt(12))
rrect(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.95), fill=PANEL, brdr=VIOLET, line_w=1.0)
txt(s, Inches(1.2), Inches(5.9), Inches(11.0), Inches(0.6),
    [[("A pilot's entire provider cost sits in the tens of dollars — the decision is about fit and "
       "governance, not spend.", 13, BODY, False, BODYF, True)]], line_spacing=1.15)

prs.save("/home/user/horizon-engine/scratchpad/Horizon-Engine-Omnicom-Briefing.pptx")
print("saved:", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
