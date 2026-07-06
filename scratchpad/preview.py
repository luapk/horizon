"""Approximate renderer: draw each slide's real shape geometry + text from the
pptx to a PNG so layout (overlap / overflow / off-slide) can be eyeballed."""
from pptx import Presentation
from pptx.util import Emu
from PIL import Image, ImageDraw, ImageFont
import textwrap

SCALE = 96 / 914400  # EMU -> px at 96dpi
W, H = int(13.333 * 96), int(7.5 * 96)

def font(sz, bold=False):
    for name in (["DejaVuSans-Bold.ttf"] if bold else ["DejaVuSans.ttf"]):
        for path in (f"/usr/share/fonts/truetype/dejavu/{name}",):
            try:
                return ImageFont.truetype(path, max(8, int(sz)))
            except Exception:
                pass
    return ImageFont.load_default()

def rgb(c):
    try:
        return (c[0], c[1], c[2])
    except Exception:
        return None

prs = Presentation("Horizon-Engine-Omnicom-Briefing.pptx")
for idx, slide in enumerate(prs.slides, 1):
    img = Image.new("RGB", (W, H), (5, 7, 15))
    d = ImageDraw.Draw(img)
    for sh in slide.shapes:
        try:
            x, y = int(sh.left * SCALE), int(sh.top * SCALE)
            w, h = int(sh.width * SCALE), int(sh.height * SCALE)
        except Exception:
            continue
        # fill + border for autoshapes
        if sh.shape_type is not None and sh.has_text_frame is False or (hasattr(sh, "fill")):
            fill = None; line = None
            try:
                if sh.fill.type == 1:
                    fill = rgb(sh.fill.fore_color.rgb)
            except Exception:
                pass
            try:
                if sh.line.width and sh.line.color and sh.line.color.type is not None:
                    line = rgb(sh.line.color.rgb)
            except Exception:
                pass
            if fill or line:
                d.rectangle([x, y, x + w, y + h], fill=fill, outline=line)
        # tables
        if sh.has_table:
            tbl = sh.table
            cy = y
            for ri, row in enumerate(tbl.rows):
                rh = int(row.height * SCALE)
                cx = x
                for ci, cell in enumerate(row.cells):
                    cw = int(tbl.columns[ci].width * SCALE)
                    cf = None
                    try:
                        cf = rgb(cell.fill.fore_color.rgb)
                    except Exception:
                        pass
                    d.rectangle([cx, cy, cx + cw, cy + rh], fill=cf, outline=(42, 51, 72))
                    t = cell.text_frame.text
                    if t:
                        d.text((cx + 6, cy + 6), t[:26], font=font(11), fill=(200, 208, 224))
                    cx += cw
                cy += rh
        # text
        if sh.has_text_frame:
            ty = y + 3
            for p in sh.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs) or p.text
                if not t:
                    continue
                sz = 14
                col = (230, 234, 242)
                if p.runs:
                    r0 = p.runs[0]
                    if r0.font.size:
                        sz = int(r0.font.size.pt)
                    if r0.font.color and r0.font.color.type is not None:
                        col = rgb(r0.font.color.rgb) or col
                fnt = font(sz * 1.05, bold=bool(p.runs and p.runs[0].font.bold))
                maxchars = max(8, int(w / (sz * 0.6)))
                for lineTxt in textwrap.wrap(t, maxchars) or [t]:
                    d.text((x + 3, ty), lineTxt, font=fnt, fill=col)
                    ty += int(sz * 1.25)
    img.save(f"prev_{idx:02d}.png")
print("rendered", len(prs.slides._sldIdLst), "slides")
